"""Generate the HERD hackathon deck as a real .pptx.

Design language mirrors the product: warm paper, ink, and one signal colour per
verdict. Numbers and labels are set in a mono face, the same convention the UI
uses, so the deck reads as an extension of the thing it is presenting.
"""
from __future__ import annotations

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# -- palette (matches web/tailwind.config.js at the shipped revision) --------
PAPER = RGBColor(0xF5, 0xF2, 0xEB)
CARD = RGBColor(0xFC, 0xFB, 0xF8)
INK = RGBColor(0x12, 0x11, 0x0C)
MUTED = RGBColor(0x57, 0x54, 0x4A)
FAINT = RGBColor(0x8B, 0x87, 0x79)
RULE = RGBColor(0xDF, 0xDB, 0xCC)
FALSE = RGBColor(0xA3, 0x20, 0x17)
MISLEAD = RGBColor(0xA5, 0x6A, 0x00)
TRUE = RGBColor(0x14, 0x62, 0x4A)
PAPERW = RGBColor(0xFD, 0xFD, 0xFC)

SANS = "Segoe UI"
SANS_L = "Segoe UI Light"
MONO = "Consolas"

W = Inches(13.333)
H = Inches(7.5)

prs = Presentation()
prs.slide_width = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]


def slide(bg=PAPER):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    r.fill.solid()
    r.fill.fore_color.rgb = bg
    r.line.fill.background()
    r.shadow.inherit = False
    # push background to back
    sp = r._element
    sp.getparent().remove(sp)
    s.shapes._spTree.insert(2, sp)
    return s


def box(s, x, y, w, h):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    return tb, tf


def para(tf, text, size, color=INK, bold=False, font=SANS, first=False,
         space_before=0, space_after=6, align=PP_ALIGN.LEFT, spacing=1.0,
         tracking=None):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    p.line_spacing = spacing
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = font
    run.font.color.rgb = color
    if tracking is not None:
        _tracking(run, tracking)
    return p, run


def _tracking(run, pts):
    run.font._rPr.set("spc", str(int(pts * 100)))


def rect(s, x, y, w, h, color, line=None, lw=1.0):
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    r.fill.solid()
    r.fill.fore_color.rgb = color
    if line is None:
        r.line.fill.background()
    else:
        r.line.color.rgb = line
        r.line.width = Pt(lw)
    r.shadow.inherit = False
    return r


def hline(s, x, y, w, color=INK, weight=1.25):
    ln = s.shapes.add_connector(2, x, y, x + w, y)
    ln.line.color.rgb = color
    ln.line.width = Pt(weight)
    ln.shadow.inherit = False
    return ln


def kicker(s, text, x=Inches(0.9), y=Inches(0.62), color=FAINT):
    _, tf = box(s, x, y, Inches(11), Inches(0.35))
    para(tf, text.upper(), 12, color, bold=True, font=MONO, first=True,
         tracking=2.0, space_after=0)


def title(s, text, x=Inches(0.9), y=Inches(1.0), w=Inches(11.5), size=40,
          color=INK):
    _, tf = box(s, x, y, w, Inches(1.4))
    para(tf, text, size, color, bold=True, first=True, spacing=0.98,
         space_after=0)


def footer(s, n):
    _, tf = box(s, Inches(0.9), Inches(7.02), Inches(9), Inches(0.3))
    para(tf, "HERD  ·  an immune system for campus misinformation", 9,
         FAINT, font=MONO, first=True, tracking=0.8, space_after=0)
    _, tf2 = box(s, Inches(12.0), Inches(7.02), Inches(0.9), Inches(0.3))
    para(tf2, f"{n:02d}", 9, FAINT, font=MONO, first=True,
         align=PP_ALIGN.RIGHT, space_after=0)


# ============================================================ 01 · TITLE
s = slide(INK)
# thin accent rule
rect(s, Inches(0.9), Inches(2.55), Inches(1.4), Pt(4), FALSE)
_, tf = box(s, Inches(0.86), Inches(2.75), Inches(11.6), Inches(2.2))
para(tf, "HERD", 128, PAPERW, bold=True, first=True, spacing=0.9,
     space_after=0)
_, tf = box(s, Inches(0.92), Inches(4.9), Inches(11), Inches(1.0))
para(tf, "An immune system for campus misinformation.", 26,
     RGBColor(0xE9, 0xE6, 0xDD), first=True, space_after=4)
para(tf, "It investigates a message the way a careful person would — cheapest "
     "checks first — and stops the moment it can honestly stop.", 15,
     RGBColor(0xB8, 0xB4, 0xA8), spacing=1.25)
_, tf = box(s, Inches(0.92), Inches(6.55), Inches(11.5), Inches(0.5))
para(tf, "ECHO 2026  ·  “BUILD BY SUNSET”  ·  VNR VJIET × STUDENTALUMNI.AI",
     12, RGBColor(0x8B, 0x87, 0x79), font=MONO, first=True, tracking=1.6)
# corner label
_, tf = box(s, Inches(10.4), Inches(0.6), Inches(2.2), Inches(0.4))
para(tf, "CLAIM · CASCADE · VERDICT", 10, RGBColor(0x8B, 0x87, 0x79),
     font=MONO, first=True, tracking=1.4, align=PP_ALIGN.RIGHT)

# ============================================================ 02 · PROBLEM
s = slide()
kicker(s, "The problem")
title(s, "The rumour always beats the correction.")
hline(s, Inches(0.9), Inches(2.05), Inches(11.53), RULE, 1.0)
_, tf = box(s, Inches(0.9), Inches(2.35), Inches(6.6), Inches(4))
para(tf, "A message lands in the class group:", 15, MUTED, first=True,
     space_after=10)
para(tf, "“Amazon off-campus drive for the 2026 batch — register here, "
     "limited slots.”", 20, INK, bold=True, spacing=1.15, space_after=14)
para(tf, "It has a logo. It has a deadline. Forty people forward it before "
     "anyone checks. Some register. Some pay the ₹750 “registration fee.”",
     15, MUTED, spacing=1.3, space_after=10)
para(tf, "Three days later someone finally says: “guys, this is fake.”  Too "
     "late.", 15, MUTED, spacing=1.3, space_after=10)
para(tf, "The same thing happens every week — exams postponed, fest "
     "cancelled, fee deadline extended, this company is hiring.", 15, MUTED,
     spacing=1.3)
# right stat panel
rect(s, Inches(8.0), Inches(2.35), Inches(4.43), Inches(3.7), CARD, RULE, 1.0)
_, tf = box(s, Inches(8.35), Inches(2.7), Inches(3.8), Inches(3.2))
para(tf, "WHY IT HURTS", 11, FAINT, bold=True, font=MONO, first=True,
     tracking=1.6, space_after=14)
for big, small in [
    ("Minutes", "for a forward to reach the whole cohort"),
    ("Days", "for a human correction to catch up"),
    ("The people who get hurt", "are the ones who never doubted it"),
]:
    para(tf, big, 22, FALSE, bold=True, space_after=1, space_before=6)
    para(tf, small, 13, MUTED, spacing=1.15)
footer(s, 2)

# ============================================================ 03 · INSIGHT
s = slide(INK)
kicker(s, "The reframe", color=FAINT)
_, tf = box(s, Inches(0.9), Inches(1.15), Inches(11.5), Inches(1.6))
para(tf, "Treat a rumour as an infection,", 44, PAPERW, bold=True, first=True,
     spacing=1.0, space_after=0)
para(tf, "not a document.", 44, FALSE, bold=True, spacing=1.0)
hline(s, Inches(0.9), Inches(3.15), Inches(11.53), RGBColor(0x3A,0x38,0x30), 1.0)
_, tf = box(s, Inches(0.9), Inches(3.45), Inches(11.4), Inches(1.4))
para(tf, "“Is this true?” is the easy question — and it always arrives too "
     "late, because the person who bothers to check was never the one at "
     "risk.", 17, RGBColor(0xCF, 0xCB, 0xC1), first=True, spacing=1.3)
_, tf = box(s, Inches(0.9), Inches(4.95), Inches(11.4), Inches(1.6))
para(tf, "The useful question:", 13, FAINT, font=MONO, first=True,
     tracking=1.2, space_after=8)
para(tf, "How fast is this spreading, who hasn’t been reached yet, "
     "and can I get there first?", 24, PAPERW, bold=True, spacing=1.15)
_, tf = box(s, Inches(0.9), Inches(6.4), Inches(11.4), Inches(0.6))
para(tf, "That is an interception problem, not a classification problem. "
     "Epidemiology, not fact-checking.", 14, RGBColor(0x9A,0x96,0x8A),
     first=True)

# ============================================================ 04 · WHAT IT DOES
s = slide()
kicker(s, "What HERD does")
title(s, "One suspicious message in.\nA cited, honest verdict out.")
hline(s, Inches(0.9), Inches(2.55), Inches(11.53), RULE, 1.0)
steps = [
    ("Perceive", "OCR + claim extraction. Turns a messy, code-mixed "
     "Telugu/Hindi/English screenshot into a structured, falsifiable claim."),
    ("Recognise", "Strain matching. Already seen — here or at any campus? "
     "Answer instantly from memory instead of re-investigating."),
    ("Investigate", "A four-tier cascade of specialist agents that return "
     "cited evidence, cheapest checks first, exiting early when it can."),
    ("Judge", "Deterministic log-odds aggregation sets the verdict and a "
     "confidence. The model only writes the prose."),
]
x = Inches(0.9)
cw = Inches(2.83)
gap = Inches(0.06)
for i, (h, b) in enumerate(steps):
    cx = x + (cw + gap) * i
    rect(s, cx, Inches(2.95), cw, Inches(3.15), CARD, RULE, 1.0)
    rect(s, cx, Inches(2.95), cw, Pt(4), INK)
    _, tf = box(s, cx + Inches(0.28), Inches(3.2), cw - Inches(0.56), Inches(2.8))
    para(tf, f"0{i+1}", 13, FAINT, bold=True, font=MONO, first=True,
         tracking=1.0, space_after=8)
    para(tf, h, 20, INK, bold=True, space_after=8)
    para(tf, b, 13, MUTED, spacing=1.28)
footer(s, 4)

# ============================================================ 05 · CASCADE
s = slide()
kicker(s, "The engine")
title(s, "A cascade that spends the least\nevidence it can get away with.")
_, tf = box(s, Inches(0.9), Inches(2.25), Inches(11.5), Inches(0.5))
para(tf, "Cheapest checks first. The cascade exits the moment a tier settles "
     "it — asymmetrically, because being wrong about a real notice costs more "
     "than one extra check.", 14, MUTED, first=True, spacing=1.2)
tiers = [
    ("T0", "Wording", "Costs nothing", "Pressure, payment shape, template "
     "lineage. FraudHeuristics · TemplateProvenance · StrainPrior", TRUE),
    ("T1", "Infrastructure", "Cheap probes", "Domains, links and contact "
     "details vs. what the campus actually publishes. DomainForensics · "
     "URLSafety · ContactForensics", TRUE),
    ("T2", "Official record", "Retrieval", "The only tier allowed to confirm "
     "a claim is genuine. InstitutionalSource · OfficialChannel", MISLEAD),
    ("T3", "Open web", "Terminal", "Bought only when everything cheaper had "
     "to abstain. OpenWebResearch", FALSE),
]
y = Inches(2.95)
for tag, name, cost, body, col in tiers:
    rect(s, Inches(0.9), y, Pt(4), Inches(0.86), col)
    _, tf = box(s, Inches(1.15), y - Inches(0.02), Inches(1.4), Inches(0.9))
    para(tf, tag, 22, col, bold=True, font=MONO, first=True, space_after=0)
    _, tf = box(s, Inches(2.35), y - Inches(0.03), Inches(3.0), Inches(0.9))
    para(tf, name, 18, INK, bold=True, first=True, space_after=1)
    para(tf, cost.upper(), 10, FAINT, font=MONO, tracking=1.2)
    _, tf = box(s, Inches(5.3), y - Inches(0.02), Inches(7.1), Inches(0.9))
    para(tf, body, 12.5, MUTED, first=True, spacing=1.12)
    y += Inches(0.98)
footer(s, 5)

# ============================================================ 06 · EVIDENCE
s = slide()
kicker(s, "The discipline")
title(s, "Agents return evidence, never verdicts.")
hline(s, Inches(0.9), Inches(2.05), Inches(11.53), RULE, 1.0)
_, tf = box(s, Inches(0.9), Inches(2.35), Inches(5.6), Inches(4.2))
para(tf, "Two contracts, enforced at the type level:", 15, INK, bold=True,
     first=True, space_after=14)
para(tf, "Cite or stay silent.", 17, INK, bold=True, space_after=3)
para(tf, "A non-neutral finding is rejected unless it carries at least one "
     "source. An uncited assertion cannot physically exist.", 13.5, MUTED,
     spacing=1.28, space_after=14)
para(tf, "TRUE requires confirmation.", 17, INK, bold=True, space_after=3)
para(tf, "Absence of fraud indicators is not evidence of authenticity — it "
     "is equally consistent with a well-made scam. Only a confirming source "
     "can say “genuine.”", 13.5, MUTED, spacing=1.28)
# right: the judgement side
rect(s, Inches(6.9), Inches(2.35), Inches(5.53), Inches(4.0), INK)
_, tf = box(s, Inches(7.25), Inches(2.7), Inches(4.85), Inches(3.4))
para(tf, "WHO DECIDES", 11, FAINT, bold=True, font=MONO, first=True,
     tracking=1.6, space_after=14)
para(tf, "The verdict is arithmetic.", 20, PAPERW, bold=True, space_after=6)
para(tf, "Evidence is combined in log-odds by a deterministic aggregator "
     "with calibrated constants. The system produces a verdict with no "
     "network and no language model at all.", 13.5,
     RGBColor(0xC6,0xC2,0xB8), spacing=1.3, space_after=12)
para(tf, "The LLM only writes the prose.", 16, FALSE, bold=True,
     space_after=4)
para(tf, "“The judgement is not generated. The sentence explaining it is.”",
     13, RGBColor(0x9A,0x96,0x8A), spacing=1.25)
footer(s, 6)

# ============================================================ 07 · BELIEF AXIS
s = slide()
kicker(s, "How a verdict is reached")
title(s, "Belief moves in log-odds, band to band.")
_, tf = box(s, Inches(0.9), Inches(2.15), Inches(11.5), Inches(0.5))
para(tf, "Every investigation is a path: a prior, then each tier nudging "
     "belief left (genuine) or right (fabricated), landing in one of four "
     "bands.", 14, MUTED, first=True, spacing=1.2)
# axis
ax_x, ax_y, ax_w = Inches(0.95), Inches(3.55), Inches(11.4)
bands = [("CONFIRMED", TRUE, 0.20), ("NO CLAIM", RGBColor(0x55,0x52,0x4A), 0.45),
         ("DISTORTED", MISLEAD, 0.25), ("FABRICATED", FALSE, 0.10)]
cx = ax_x
for name, col, frac in bands:
    bw = Emu(int(ax_w * frac))
    rect(s, cx, ax_y, bw - Pt(3), Inches(0.5), col)
    _, tf = box(s, cx, ax_y + Inches(0.62), bw, Inches(0.4))
    para(tf, name, 11, col, bold=True, font=MONO, first=True, tracking=1.0,
         align=PP_ALIGN.CENTER)
    cx = Emu(int(cx) + int(bw))
# prior + landing markers
prior_x = ax_x + Emu(int(ax_w * 0.35))
land_x = ax_x + Emu(int(ax_w * 0.955))
rect(s, prior_x, ax_y - Inches(0.35), Pt(2), Inches(1.2), FAINT)
_, tf = box(s, prior_x - Inches(0.6), ax_y - Inches(0.7), Inches(1.2), Inches(0.35))
para(tf, "PRIOR 0.35", 10, FAINT, font=MONO, first=True, align=PP_ALIGN.CENTER)
rect(s, land_x, ax_y - Inches(0.35), Pt(4), Inches(1.2), FALSE)
_, tf = box(s, land_x - Inches(1.0), ax_y - Inches(0.7), Inches(1.4), Inches(0.35))
para(tf, "0.998", 12, FALSE, bold=True, font=MONO, first=True,
     align=PP_ALIGN.CENTER)
_, tf = box(s, Inches(0.9), Inches(5.2), Inches(11.5), Inches(1.4))
para(tf, "Why log-odds?", 15, INK, bold=True, first=True, space_after=6)
para(tf, "On a linear probability axis, 0.94 → 0.998 looks like a rounding "
     "error. In log-odds — the scale the aggregator actually adds evidence "
     "on — it is more evidence than everything before it combined. The UI "
     "plots the real thing, so the judgement reads as a judgement, not a "
     "number.", 14, MUTED, spacing=1.3)
footer(s, 7)

# ============================================================ 08 · COMPOUNDS
s = slide()
kicker(s, "Why it compounds")
title(s, "The wider an attack spreads, the cheaper\nit becomes to neutralise.")
hline(s, Inches(0.9), Inches(2.55), Inches(11.53), RULE, 1.0)
_, tf = box(s, Inches(0.9), Inches(2.85), Inches(11.4), Inches(1.1))
para(tf, "Scams scale by repetition — the same template, resent forever. "
     "HERD turns that strength into a weakness. Report #1 costs a full "
     "investigation; reports #2–#4000 cost a lookup. One student’s "
     "investigation becomes permanent immunity for everyone who follows.",
     15, MUTED, first=True, spacing=1.3)
cols = [
    ("Strain memory is global", "A template that cost one campus a full "
     "investigation is recognised instantly at the next one.", INK),
    ("Evidence is scoped", "Each campus’s verdict is still derived from its "
     "own notice board, its own official channels. Share the pattern, scope "
     "the proof.", INK),
    ("Cold-start inverts", "HERD is most valuable to the newest institution — "
     "it arrives carrying everyone else’s accumulated immunity.", FALSE),
]
x = Inches(0.9)
cw = Inches(3.71)
for i, (h, b, col) in enumerate(cols):
    cx = x + (cw + Inches(0.2)) * i
    rect(s, cx, Inches(4.25), cw, Inches(2.3), CARD, RULE, 1.0)
    _, tf = box(s, cx + Inches(0.3), Inches(4.5), cw - Inches(0.6), Inches(1.9))
    para(tf, h, 17, col, bold=True, first=True, space_after=8, spacing=1.05)
    para(tf, b, 13, MUTED, spacing=1.28)
footer(s, 8)

# ============================================================ 09 · DEMO / RESULT
s = slide(INK)
kicker(s, "Live investigation", color=FAINT)
title(s, "A real scam, judged in ~1.2 seconds.", color=PAPERW)
_, tf = box(s, Inches(0.9), Inches(2.0), Inches(11.4), Inches(0.9))
para(tf, "“URGENT!! Placement drive registration closing TODAY. Pay ₹5000 to "
     "secure your slot… Register at vnrvjiet-placements.online. Forward to "
     "all your groups.”", 14, RGBColor(0xC6,0xC2,0xB8), first=True,
     spacing=1.25)
# verdict block
rect(s, Inches(0.9), Inches(3.15), Inches(5.4), Inches(3.4),
     RGBColor(0x1B,0x1A,0x14))
rect(s, Inches(0.9), Inches(3.15), Inches(5.4), Pt(5), FALSE)
_, tf = box(s, Inches(1.2), Inches(3.4), Inches(4.9), Inches(3.0))
para(tf, "VERDICT", 11, FAINT, bold=True, font=MONO, first=True, tracking=1.6,
     space_after=6)
para(tf, "Fabricated", 46, FALSE, bold=True, space_after=6)
para(tf, "FALSE · posterior 0.998 · confidence 0.99", 13,
     RGBColor(0xC6,0xC2,0xB8), font=MONO, space_after=14)
para(tf, "4 tiers ran · 0 left unbought · all 9 agents returned evidence",
     12.5, RGBColor(0x9A,0x96,0x8A), spacing=1.2)
# evidence list
rect(s, Inches(6.55), Inches(3.15), Inches(5.88), Inches(3.4),
     RGBColor(0x1B,0x1A,0x14))
_, tf = box(s, Inches(6.85), Inches(3.4), Inches(5.3), Inches(3.0))
para(tf, "WHAT SPOKE", 11, FAINT, bold=True, font=MONO, first=True,
     tracking=1.6, space_after=10)
for name, find in [
    ("FraudHeuristics", "upfront fee · personal UPI · <48h deadline"),
    ("DomainForensics", "look-alike domain, not a campus domain"),
    ("InstitutionalSource", "no matching notice on any official source"),
    ("OfficialChannel", "no matching broadcast on any channel"),
]:
    p = tf.add_paragraph()
    p.space_after = Pt(7)
    p.line_spacing = 1.05
    r1 = p.add_run(); r1.text = name + "  "
    r1.font.size = Pt(13); r1.font.bold = True; r1.font.name = MONO
    r1.font.color.rgb = PAPERW
    r2 = p.add_run(); r2.text = find
    r2.font.size = Pt(12); r2.font.name = SANS
    r2.font.color.rgb = RGBColor(0xB0,0xAC,0xA0)

# ============================================================ 10 · METRICS
s = slide()
kicker(s, "Where it stands")
title(s, "Measured, not estimated.")
hline(s, Inches(0.9), Inches(1.95), Inches(11.53), RULE, 1.0)
metrics = [
    ("198", "tests passing", "0 failing · ~35s", TRUE),
    ("9 / 9", "agents live", "returning real evidence", TRUE),
    ("13 / 18", "scams caught", "as FALSE", FALSE),
    ("0", "genuine notices libelled", "safety first", TRUE),
    ("~1.2s", "warm investigation", "all four tiers", INK),
    ("38%", "exact-label match", "100% correct-or-abstain", MISLEAD),
]
gx, gy = Inches(0.9), Inches(2.35)
cw, ch = Inches(3.71), Inches(1.95)
for i, (big, lab, sub, col) in enumerate(metrics):
    r, c = divmod(i, 3)
    cx = gx + (cw + Inches(0.2)) * c
    cy = gy + (ch + Inches(0.2)) * r
    rect(s, cx, cy, cw, ch, CARD, RULE, 1.0)
    _, tf = box(s, cx + Inches(0.3), cy + Inches(0.22), cw - Inches(0.6), ch - Inches(0.4))
    para(tf, big, 40, col, bold=True, font=MONO, first=True, space_after=2)
    para(tf, lab, 15, INK, bold=True, space_after=1)
    para(tf, sub, 12, FAINT, spacing=1.1)
_, tf = box(s, Inches(0.9), Inches(6.5), Inches(11.5), Inches(0.5))
para(tf, "Read the 38% correctly: every remaining case is an honest "
     "abstention, not a harmful mistake. Zero genuine notices are ever "
     "accused.", 12.5, MUTED, first=True, spacing=1.15)
footer(s, 10)

# ============================================================ 11 · RIGOR
s = slide()
kicker(s, "Engineering rigor")
title(s, "The parts a demo doesn’t show.")
hline(s, Inches(0.9), Inches(2.05), Inches(11.53), RULE, 1.0)
items = [
    ("No numeric literal in app/", "Every tunable lives in a config file; an "
     "AST test fails the build if a constant is hard-coded — so the "
     "calibrated numbers are genuinely derived, not hand-tuned in place."),
    ("The calibrator is honest", "It ranks settings by worst-case safety "
     "margin, reports its trade-offs, and re-derives the constants from the "
     "corpus. We fixed its objective, not the number it emits."),
    ("Institution-neutral by law", "A linter forbids any campus string in "
     "app/ or web/. A new campus is one YAML file — zero code changes."),
    ("29 architecture decisions", "Every significant fork is recorded with "
     "the options considered and the consequences accepted. Three were "
     "written mid-build, because measurement contradicted the design."),
]
x = Inches(0.9)
cw = Inches(5.66)
for i, (h, b) in enumerate(items):
    r, c = divmod(i, 2)
    cx = x + (cw + Inches(0.2)) * c
    cy = Inches(2.4) + (Inches(2.05) + Inches(0.2)) * r
    rect(s, cx, cy, Pt(4), Inches(1.9), INK)
    _, tf = box(s, cx + Inches(0.28), cy, cw - Inches(0.4), Inches(2.0))
    para(tf, h, 18, INK, bold=True, first=True, space_after=7, spacing=1.05)
    para(tf, b, 13, MUTED, spacing=1.28)
footer(s, 11)

# ============================================================ 12 · PRIVACY
s = slide()
kicker(s, "Trust & safety")
title(s, "Report-driven. Consent-based.\nIt never watches anyone.")
hline(s, Inches(0.9), Inches(2.75), Inches(11.53), RULE, 1.0)
_, tf = box(s, Inches(0.9), Inches(3.1), Inches(11.4), Inches(1.2))
para(tf, "HERD never reads a group, never joins a chat, never monitors "
     "anyone. It only ever sees what a human explicitly hands it — "
     "intercepting a habit that already exists, since people already forward "
     "a suspicious message to a friend asking “is this real?”", 16, MUTED,
     first=True, spacing=1.35)
chips = [
    ("Redacted on arrival", "Phone numbers and UPI handles are masked before "
     "anything is stored — but the handle’s provider survives, because it is "
     "evidence."),
    ("Pseudonymous by salt", "Reporters are a rotating hash, never an "
     "identity. There is nothing to de-anonymise."),
    ("Offline-first", "The verdict is produced with no network and no LLM. "
     "Nothing is required to leave the machine."),
]
x = Inches(0.9)
cw = Inches(3.71)
for i, (h, b) in enumerate(chips):
    cx = x + (cw + Inches(0.2)) * i
    rect(s, cx, Inches(4.5), cw, Inches(2.05), CARD, RULE, 1.0)
    _, tf = box(s, cx + Inches(0.3), Inches(4.72), cw - Inches(0.6), Inches(1.7))
    para(tf, h, 16, TRUE, bold=True, first=True, space_after=7, spacing=1.05)
    para(tf, b, 12.5, MUTED, spacing=1.25)
footer(s, 12)

# ============================================================ 13 · WHAT'S NEXT
s = slide()
kicker(s, "Honest limitations & next")
title(s, "What we’d build after sunset.")
hline(s, Inches(0.9), Inches(2.05), Inches(11.53), RULE, 1.0)
rows = [
    ("Confirmation coverage", "Tier 2 is what turns abstentions into "
     "confirmed TRUE. More official sources per campus directly lifts the "
     "exact-match rate — the current gap is missing confirmers, not wrong "
     "verdicts."),
    ("Pre-bunk delivery loop", "The inoculation card is generated; wiring it "
     "to a real admin broadcast closes the interception loop the whole thesis "
     "rests on."),
    ("Cassette-backed demo", "Record/replay of the network tier so the live "
     "run is deterministic and offline, immune to rate limits on stage."),
    ("Cross-campus network", "Onboard a second institution and demonstrate a "
     "strain caught at campus B the instant it was investigated at campus A."),
]
y = Inches(2.35)
for h, b in rows:
    rect(s, Inches(0.9), y + Inches(0.04), Inches(0.14), Inches(0.14), FALSE)
    _, tf = box(s, Inches(1.25), y - Inches(0.05), Inches(11.1), Inches(1.0))
    p = tf.paragraphs[0]
    p.line_spacing = 1.2
    r1 = p.add_run(); r1.text = h + " — "
    r1.font.size = Pt(15); r1.font.bold = True; r1.font.name = SANS
    r1.font.color.rgb = INK
    r2 = p.add_run(); r2.text = b
    r2.font.size = Pt(13.5); r2.font.name = SANS; r2.font.color.rgb = MUTED
    y += Inches(1.08)
footer(s, 13)

# ============================================================ 14 · CLOSING
s = slide(INK)
rect(s, Inches(0.9), Inches(2.3), Inches(1.4), Pt(4), FALSE)
_, tf = box(s, Inches(0.86), Inches(2.5), Inches(11.6), Inches(1.6))
para(tf, "HERD", 96, PAPERW, bold=True, first=True, spacing=0.9,
     space_after=0)
_, tf = box(s, Inches(0.92), Inches(4.25), Inches(11.4), Inches(1.4))
para(tf, "Not a fact-checker. An immune system.", 26, PAPERW, bold=True,
     first=True, space_after=8)
para(tf, "It reads the message, spends the least evidence it can, cites "
     "everything, and tells you when it doesn’t know.", 15,
     RGBColor(0xB8,0xB4,0xA8), spacing=1.3)
hline(s, Inches(0.92), Inches(6.1), Inches(11.5), RGBColor(0x3A,0x38,0x30), 1.0)
_, tf = box(s, Inches(0.92), Inches(6.3), Inches(11.5), Inches(0.6))
para(tf, "github.com/thirumani-vihaan/hackathon   ·   198 tests · 9/9 agents · "
     "0 harmful   ·   Thank you.", 13, RGBColor(0x9A,0x96,0x8A), font=MONO,
     first=True, tracking=0.6)

out = "HERD-hackathon-deck.pptx"
prs.save(out)
print(f"saved {out} with {len(prs.slides.__iter__.__self__._sldIdLst)} slides")
